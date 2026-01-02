from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

def create_basic_presentation(title, output_file):
    # Create a presentation object
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title
    subtitle.text = f"Created on {datetime.now().strftime('%Y-%m-%d')}"
    
    # Add a slide with title and content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'First Content Slide'
    
    tf = body_shape.text_frame
    tf.text = 'This is the first content slide.'
    
    # Add a bullet point slide
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Bullet Points'
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    tf.text = 'First Level'
    p = tf.add_paragraph()
    p.text = 'Second Level'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Another First Level'
    
    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == "__main__":
    presentation_title = "My Presentation"
    output_filename = "presentation.pptx"
    create_basic_presentation(presentation_title, output_filename)
